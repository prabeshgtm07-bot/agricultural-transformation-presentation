"""
Agricultural Transformation - PowerPoint Generator
Comprehensive 22-slide presentation with professional green/earth-tone design (16:9)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# ── Palette (earth tones / greens) ────────────────────────────────────────────
C_DARK_GREEN   = RGBColor(27,  94,  32)   # deep forest green  #1B5E20
C_MID_GREEN    = RGBColor(46, 125,  50)   # mid green          #2E7D32
C_LIGHT_GREEN  = RGBColor(200, 230, 201)  # pale mint          #C8E6C9
C_EARTH_BROWN  = RGBColor(93,  64,  55)   # earthy brown       #5D4037
C_WARM_AMBER   = RGBColor(230, 162,  60)  # warm amber/gold    #E6A23C
C_CREAM        = RGBColor(250, 245, 235)  # off-white cream    #FAF5EB
C_WHITE        = RGBColor(255, 255, 255)
C_DARK_GRAY    = RGBColor(33,  33,  33)
C_MID_GRAY     = RGBColor(97,  97,  97)

# ── Slide dimensions  16:9  ───────────────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # truly blank layout

# ══════════════════════════════════════════════════════════════════════════════
# Helper utilities
# ══════════════════════════════════════════════════════════════════════════════

def _solid_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def _textbox(slide, left, top, width, height, text, font_size,
             bold=False, color=None, align=PP_ALIGN.LEFT,
             italic=False, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def _add_speaker_note(slide, note_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = note_text


# ══════════════════════════════════════════════════════════════════════════════
# Section-header accent bar (left vertical stripe)
# ══════════════════════════════════════════════════════════════════════════════

def _accent_bar(slide):
    _rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, C_DARK_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# Slide builders
# ══════════════════════════════════════════════════════════════════════════════

def add_title_slide():
    """Slide 1 – Title"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_DARK_GREEN)

    # decorative bottom band
    _rect(slide, Inches(0), Inches(6.3), SLIDE_W, Inches(1.2), C_MID_GREEN)

    # amber accent stripe
    _rect(slide, Inches(0), Inches(5.95), SLIDE_W, Inches(0.15), C_WARM_AMBER)

    # main title
    _textbox(slide,
             Inches(1.0), Inches(1.2), Inches(11.33), Inches(2.0),
             "Agricultural Transformation",
             font_size=54, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # subtitle
    _textbox(slide,
             Inches(1.0), Inches(3.1), Inches(11.33), Inches(0.9),
             "Driving Sustainable Growth Through Technology, Innovation & Policy",
             font_size=22, color=C_WARM_AMBER, align=PP_ALIGN.CENTER)

    # meta line
    _textbox(slide,
             Inches(1.0), Inches(4.1), Inches(11.33), Inches(0.55),
             "April 2026  |  Global Agriculture Forum",
             font_size=16, color=C_LIGHT_GREEN, align=PP_ALIGN.CENTER)

    # presenter
    _textbox(slide,
             Inches(1.0), Inches(6.35), Inches(11.33), Inches(0.85),
             "Presented by: Agricultural Transformation Research Team",
             font_size=14, color=C_WHITE, align=PP_ALIGN.CENTER)

    _add_speaker_note(slide,
        "Welcome the audience. This presentation covers the global imperative for agricultural "
        "transformation, spanning technology adoption, sustainable practices, case studies, "
        "and a clear call to action. Estimated delivery time: 35-40 minutes.")


def add_toc_slide():
    """Slide 2 – Table of Contents"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _accent_bar(slide)

    _textbox(slide, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.7),
             "Table of Contents", font_size=34, bold=True, color=C_DARK_GREEN)

    _rect(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.05), C_WARM_AMBER)

    sections = [
        ("01", "Introduction – What Is Agricultural Transformation?",   "Slides 3–5"),
        ("02", "Key Areas of Transformation",                           "Slides 6–10"),
        ("03", "Case Studies – Real-World Success",                     "Slides 11–13"),
        ("04", "Benefits & Impact",                                     "Slides 14–16"),
        ("05", "Challenges & Solutions",                                "Slides 17–18"),
        ("06", "Future Outlook & Emerging Trends",                      "Slides 19–20"),
        ("07", "Call to Action",                                        "Slide 21"),
        ("08", "Conclusion & Q&A",                                      "Slide 22"),
    ]

    row_h = Inches(0.62)
    for i, (num, label, pages) in enumerate(sections):
        y = Inches(1.1) + i * row_h
        bg = C_LIGHT_GREEN if i % 2 == 0 else C_WHITE
        _rect(slide, Inches(0.35), y, Inches(12.6), row_h - Inches(0.05), bg)

        # number badge
        badge = _rect(slide, Inches(0.4), y + Inches(0.07),
                      Inches(0.45), Inches(0.45), C_MID_GREEN)

        _textbox(slide, Inches(0.4), y + Inches(0.07), Inches(0.45), Inches(0.45),
                 num, font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        _textbox(slide, Inches(1.0), y + Inches(0.1), Inches(10.3), Inches(0.42),
                 label, font_size=15, color=C_DARK_GRAY)

        _textbox(slide, Inches(11.5), y + Inches(0.1), Inches(1.4), Inches(0.42),
                 pages, font_size=13, color=C_MID_GRAY, align=PP_ALIGN.RIGHT)

    _add_speaker_note(slide,
        "Walk through each section briefly. This agenda gives the audience a clear roadmap "
        "for the presentation and sets expectations for the depth of coverage.")


def _content_header(slide, title, section_label):
    """Reusable header: green title bar + section tag"""
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.05), C_DARK_GREEN)
    _textbox(slide, Inches(0.35), Inches(0.12), Inches(11.0), Inches(0.75),
             title, font_size=30, bold=True, color=C_WHITE)
    _textbox(slide, Inches(11.4), Inches(0.25), Inches(1.8), Inches(0.5),
             section_label, font_size=11, color=C_WARM_AMBER, align=PP_ALIGN.RIGHT)


def _bullet_list(slide, items, left, top, width, height,
                 font_size=16, bullet="▸ ", color=None):
    if color is None:
        color = C_DARK_GRAY
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet + item if item.strip() else ""
        p.space_before = Pt(4)
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


# ─────────────────────────────────────────────
# INTRODUCTION SECTION  (Slides 3-5)
# ─────────────────────────────────────────────

def add_intro_what():
    """Slide 3 – What is agricultural transformation?"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _content_header(slide, "What Is Agricultural Transformation?", "Introduction")

    _textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.55),
             "A systemic shift from subsistence farming to modern, productive, and sustainable food systems.",
             font_size=17, italic=True, color=C_EARTH_BROWN)

    bullets = [
        "Modernizing farming methods, tools, and market linkages",
        "Shifting from low-input subsistence to high-value commercial production",
        "Integrating technology, data, and innovation across the agricultural value chain",
        "Empowering smallholder farmers through knowledge, finance, and infrastructure",
        "Building resilient food systems that can withstand climate and market shocks",
        "Aligning rural economies with national development and food security goals",
    ]
    _bullet_list(slide, bullets, Inches(0.5), Inches(1.85), Inches(12.3), Inches(5.0))

    _add_speaker_note(slide,
        "Agricultural transformation is not just about new technology – it is a holistic "
        "restructuring of how food is grown, processed, and delivered. It involves farmers, "
        "governments, private sector, and consumers working in an interconnected system. "
        "Emphasize that this is a proven pathway: every high-income country went through it.")


def add_intro_why():
    """Slide 4 – Why it matters"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _content_header(slide, "Why Agricultural Transformation Matters", "Introduction")

    # Two column layout
    col_w = Inches(5.9)
    col_gap = Inches(0.3)
    left1 = Inches(0.4)
    left2 = left1 + col_w + col_gap

    _rect(slide, left1, Inches(1.15), col_w, Inches(5.8), C_LIGHT_GREEN)
    _rect(slide, left2, Inches(1.15), col_w, Inches(5.8), RGBColor(255, 243, 224))

    _textbox(slide, left1 + Inches(0.15), Inches(1.2), col_w - Inches(0.3), Inches(0.5),
             "Global Imperative", font_size=17, bold=True, color=C_DARK_GREEN)

    imp_bullets = [
        "9.7 billion people to feed by 2050",
        "Agriculture employs 26% of the global workforce",
        "40% of global land is agricultural land",
        "Current yield growth: 1.6% vs needed 2.4% annually",
        "Food systems generate 30% of global GHG emissions",
        "820 million people remain food insecure today",
    ]
    _bullet_list(slide, imp_bullets, left1 + Inches(0.15), Inches(1.75),
                 col_w - Inches(0.3), Inches(4.8), font_size=15, color=C_DARK_GRAY)

    _textbox(slide, left2 + Inches(0.15), Inches(1.2), col_w - Inches(0.3), Inches(0.5),
             "Economic Opportunity", font_size=17, bold=True, color=C_EARTH_BROWN)

    opp_bullets = [
        "$8.4 trillion global agri-food market by 2030",
        "Digital agriculture market growing at 12% CAGR",
        "Precision farming saves 15–25% on input costs",
        "Value-chain modernization triples farmer income",
        "Agri-tech investment reached $51 billion in 2023",
        "Every $1 in ag R&D yields $20 in economic return",
    ]
    _bullet_list(slide, opp_bullets, left2 + Inches(0.15), Inches(1.75),
                 col_w - Inches(0.3), Inches(4.8), font_size=15, color=C_DARK_GRAY)

    _add_speaker_note(slide,
        "This slide contrasts the imperative with the opportunity. Use these statistics to "
        "create urgency. The key message is: inaction is expensive; transformation pays dividends. "
        "Agriculture is simultaneously the world's biggest challenge and biggest opportunity.")


def add_intro_challenges():
    """Slide 5 – Current global challenges"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _content_header(slide, "Current Global Challenges in Agriculture", "Introduction")

    challenges = [
        ("🌡  Climate Change",
         "Erratic rainfall, droughts, floods – crop failures rising 25% per decade"),
        ("🌱  Soil Degradation",
         "33% of global soils degraded; $40B annual productivity loss"),
        ("💧  Water Scarcity",
         "Agriculture uses 70% of freshwater; aquifers depleting in key growing regions"),
        ("👨‍🌾  Aging Farmer Population",
         "Average farmer age: 60 in developed, 45 in developing nations; youth exodus"),
        ("📦  Supply Chain Inefficiency",
         "30–40% post-harvest losses; price volatility hurts farmers & consumers"),
        ("💰  Finance Gap",
         "$170 billion annual financing gap for smallholder farmers globally"),
    ]

    card_w = Inches(3.9)
    card_h = Inches(1.45)
    cols = [Inches(0.35), Inches(4.65), Inches(8.95)]
    rows = [Inches(1.2), Inches(2.8)]
    colors = [C_LIGHT_GREEN, RGBColor(255, 243, 224),
              RGBColor(232, 245, 233), RGBColor(255, 235, 210),
              C_LIGHT_GREEN, RGBColor(255, 243, 224)]

    for idx, (heading, body) in enumerate(challenges):
        col = idx % 3
        row = idx // 3
        x = cols[col]
        y = rows[row]
        _rect(slide, x, y, card_w, card_h, colors[idx])
        _textbox(slide, x + Inches(0.12), y + Inches(0.08),
                 card_w - Inches(0.2), Inches(0.42),
                 heading, font_size=14, bold=True, color=C_DARK_GREEN)
        _textbox(slide, x + Inches(0.12), y + Inches(0.5),
                 card_w - Inches(0.2), Inches(0.85),
                 body, font_size=12, color=C_DARK_GRAY)

    _add_speaker_note(slide,
        "Present each challenge with a real statistic. These six challenges are interconnected – "
        "soil degradation worsens under climate stress, water scarcity reduces yields, and the "
        "finance gap prevents farmers from investing in solutions. Transformation must address all "
        "these simultaneously.")


# ─────────────────────────────────────────────
# KEY AREAS  (Slides 6-10)
# ─────────────────────────────────────────────

def add_tech_adoption():
    """Slide 6 – Technology adoption"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _content_header(slide, "Technology Adoption: IoT, Drones & Sensors", "Key Areas of Transformation")

    techs = [
        ("IoT & Smart Sensors",
         ["Soil moisture, pH, and nutrient monitors",
          "Automated irrigation triggers – 30% water savings",
          "Real-time microclimate weather stations"]),
        ("Drone Technology",
         ["Crop health mapping via multispectral imaging",
          "Precision pesticide & fertiliser application (50% less chemicals)",
          "Seed planting in hard-to-reach terrain"]),
        ("Satellite & Remote Sensing",
         ["NDVI mapping for large-scale crop monitoring",
          "Early drought/flood prediction (2-4 weeks ahead)",
          "Yield forecasting with ±8% accuracy"]),
        ("Farm Automation",
         ["GPS-guided tractors & autonomous harvesters",
          "Robotic weeding reduces labour costs 40%",
          "Automated greenhouse management systems"]),
    ]

    card_w = Inches(6.0)
    card_h = Inches(2.35)
    positions = [
        (Inches(0.3),  Inches(1.15)),
        (Inches(6.7),  Inches(1.15)),
        (Inches(0.3),  Inches(3.7)),
        (Inches(6.7),  Inches(3.7)),
    ]
    bg_colors = [C_LIGHT_GREEN, RGBColor(255, 243, 224),
                 RGBColor(255, 243, 224), C_LIGHT_GREEN]

    for (x, y), bg, (heading, items) in zip(positions, bg_colors, techs):
        _rect(slide, x, y, card_w, card_h, bg)
        _textbox(slide, x + Inches(0.15), y + Inches(0.1),
                 card_w - Inches(0.3), Inches(0.42),
                 heading, font_size=15, bold=True, color=C_DARK_GREEN)
        _bullet_list(slide, items, x + Inches(0.15), y + Inches(0.55),
                     card_w - Inches(0.3), Inches(1.7), font_size=13)

    _add_speaker_note(slide,
        "Technology adoption is the enabler that multiplies the impact of every other "
        "transformation pillar. Drones in Kenya cut input costs by half. IoT irrigation in India "
        "improved water use efficiency by 32%. Satellite data helped Ethiopia identify crop "
        "stress areas 3 weeks before visible signs. Cost of these technologies is dropping 20% "
        "per year, making them increasingly accessible to smallholders.")


def add_sustainable_farming():
    """Slide 7 – Sustainable farming practices"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _content_header(slide, "Sustainable Farming Practices", "Key Areas of Transformation")

    _textbox(slide, Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.5),
             "Practices that increase productivity while preserving natural resources for future generations",
             font_size=16, italic=True, color=C_EARTH_BROWN)

    practices = [
        ("Conservation Agriculture",
         "Minimal tillage + crop residue retention reduces erosion 80%, sequesters carbon"),
        ("Integrated Pest Management",
         "Biological controls + targeted chemicals cut pesticide use 50% with equal protection"),
        ("Agroforestry",
         "Trees on farms improve soil, diversify income and sequester 1–4 tonne CO₂/ha/year"),
        ("Crop Rotation & Diversification",
         "Breaking pest cycles, improving soil biology, and reducing market risk"),
        ("Water-Efficient Irrigation",
         "Drip/sprinkler systems reduce water use 40–60% versus flood irrigation"),
        ("Organic & Regenerative Methods",
         "Rebuilds soil organic matter 0.5–1% per year; premium markets offer 20–30% price uplift"),
    ]

    for i, (title, body) in enumerate(practices):
        col = i % 2
        row = i // 2
        x = Inches(0.4) + col * Inches(6.55)
        y = Inches(1.75) + row * Inches(1.55)
        _rect(slide, x, y, Inches(6.2), Inches(1.45),
              C_LIGHT_GREEN if col == 0 else RGBColor(255, 243, 224))
        _textbox(slide, x + Inches(0.15), y + Inches(0.08),
                 Inches(5.9), Inches(0.42),
                 title, font_size=14, bold=True, color=C_DARK_GREEN)
        _textbox(slide, x + Inches(0.15), y + Inches(0.52),
                 Inches(5.9), Inches(0.85),
                 body, font_size=12, color=C_DARK_GRAY)

    _add_speaker_note(slide,
        "Sustainable farming is not a trade-off with productivity – it is the route to long-term "
        "productivity. Conservation agriculture in Brazil's Cerrado region increased yields while "
        "cutting fuel costs by 60%. Agroforestry in Mali reversed desertification on 5 million ha. "
        "These practices must become the default, not the exception.")


def add_climate_smart():
    """Slide 8 – Climate-smart agriculture"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _content_header(slide, "Climate-Smart Agriculture (CSA)", "Key Areas of Transformation")

    # Three-pillar CSA visual
    pillar_data = [
        (C_MID_GREEN,   "PRODUCTIVITY",
         ["Climate-resilient crop varieties",
          "Drought/flood-tolerant seeds",
          "Adaptive cropping calendars",
          "Index-based crop insurance"]),
        (C_EARTH_BROWN, "ADAPTATION",
         ["Weather-based advisory services",
          "Diversified farming systems",
          "Soil water conservation",
          "Early warning systems"]),
        (C_WARM_AMBER,  "MITIGATION",
         ["Reduced tillage & residue mgmt.",
          "Improved livestock feed efficiency",
          "Biogas from agricultural waste",
          "Agroforestry carbon sequestration"]),
    ]

    p_w = Inches(3.9)
    for i, (color, heading, items) in enumerate(pillar_data):
        x = Inches(0.35) + i * (p_w + Inches(0.2))
        _rect(slide, x, Inches(1.15), p_w, Inches(0.6), color)
        _textbox(slide, x, Inches(1.15), p_w, Inches(0.6),
                 heading, font_size=17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _rect(slide, x, Inches(1.75), p_w, Inches(5.1), C_LIGHT_GREEN)
        _bullet_list(slide, items, x + Inches(0.15), Inches(1.85),
                     p_w - Inches(0.3), Inches(4.8), font_size=14)

    _add_speaker_note(slide,
        "CSA has three mutually reinforcing pillars – the triple win. The FAO estimates CSA "
        "practices could generate $1.8 trillion in benefits by 2030. Key success story: "
        "Vietnam's 'One Must Do, Five Reductions' rice programme cut GHG by 30% while raising "
        "yields 15%. Climate finance (Green Climate Fund, CGIAR) provides growing external "
        "resources for CSA implementation.")


def add_digital_tools():
    """Slide 9 – Digital tools and data analytics"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _content_header(slide, "Digital Tools & Data Analytics", "Key Areas of Transformation")

    items = [
        ("Farm Management Software",
         "End-to-end record keeping, planning, and cost tracking on mobile/tablet platforms"),
        ("AI & Machine Learning",
         "Yield prediction models (±5% accuracy), disease detection from images in seconds"),
        ("Big Data & Market Intelligence",
         "Price forecasting, demand signals, and export-market matching for farmer groups"),
        ("Blockchain for Traceability",
         "Farm-to-fork transparency; premium price uplift of 10–25% for verified sustainable produce"),
        ("Mobile Financial Services",
         "Digital credit scoring, mobile insurance, instant payments – reaching unbanked farmers"),
        ("e-Extension Platforms",
         "Video, voice, and chatbot advisory reaching 10x more farmers per extension officer"),
    ]

    for i, (title, body) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.35) + col * Inches(6.55)
        y = Inches(1.2) + row * Inches(1.8)
        _rect(slide, x, y, Inches(6.2), Inches(1.7),
              C_LIGHT_GREEN if col == 0 else RGBColor(255, 243, 224))
        _textbox(slide, x + Inches(0.15), y + Inches(0.08),
                 Inches(5.9), Inches(0.42),
                 title, font_size=14, bold=True, color=C_DARK_GREEN)
        _textbox(slide, x + Inches(0.15), y + Inches(0.52),
                 Inches(5.9), Inches(1.1),
                 body, font_size=12, color=C_DARK_GRAY)

    _add_speaker_note(slide,
        "Digital tools are the connective tissue of modern agriculture. In India, the eNAM "
        "electronic trading platform connected 1,000+ markets, reducing price volatility for "
        "15 million farmers. In Kenya, M-Pesa-linked digital credit reached 500,000 smallholders "
        "previously excluded from formal banking. AI-based disease detection (e.g., PlantVillage) "
        "now serves 10 million farmers in sub-Saharan Africa.")


def add_supply_chain():
    """Slide 10 – Supply chain modernization"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _content_header(slide, "Supply Chain Modernization", "Key Areas of Transformation")

    _textbox(slide, Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.45),
             "Connecting farmers to markets efficiently – reducing losses, adding value, and improving farmer income share",
             font_size=15, italic=True, color=C_EARTH_BROWN)

    steps = [
        ("FARM\nPRODUCTION", "GAP-certified practices\nTraceability codes\nQuality grading"),
        ("AGGREGATION\n& STORAGE", "Rural collection centres\nCold chain logistics\nQuality sorting"),
        ("PROCESSING\n& PACKAGING", "Value-added products\nBranding & labelling\nFood safety compliance"),
        ("DISTRIBUTION\n& RETAIL", "Regional distribution hubs\nE-commerce channels\nDirect-to-consumer"),
        ("EXPORT\n& MARKETS", "SPS compliance\nBuyer relationships\nPremium market access"),
    ]

    arrow_w = Inches(0.4)
    box_w = Inches(2.2)
    box_h = Inches(3.5)
    start_x = Inches(0.3)
    y = Inches(1.8)
    colors_steps = [C_DARK_GREEN, C_MID_GREEN, C_EARTH_BROWN, C_WARM_AMBER, C_MID_GREEN]

    for i, (step_title, step_body) in enumerate(steps):
        x = start_x + i * (box_w + arrow_w)
        _rect(slide, x, y, box_w, Inches(0.55), colors_steps[i])
        _textbox(slide, x, y, box_w, Inches(0.55),
                 step_title, font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _rect(slide, x, y + Inches(0.55), box_w, box_h - Inches(0.55), C_LIGHT_GREEN)
        _textbox(slide, x + Inches(0.1), y + Inches(0.65),
                 box_w - Inches(0.2), box_h - Inches(0.75),
                 step_body, font_size=12, color=C_DARK_GRAY)
        # arrow
        if i < len(steps) - 1:
            ax = x + box_w + Inches(0.05)
            _textbox(slide, ax, y + Inches(1.5), arrow_w - Inches(0.1), Inches(0.5),
                     "►", font_size=18, color=C_DARK_GREEN, align=PP_ALIGN.CENTER)

    # farmer income share bar chart
    chart_data = ChartData()
    chart_data.categories = ["Traditional Chain", "Modernized Chain"]
    chart_data.add_series("Farmer's Income Share (%)", (32, 58))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.3), Inches(5.5), Inches(12.7), Inches(1.7),
        chart_data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Farmer Income Share: Traditional vs Modernized Chain"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = C_MID_GREEN

    _add_speaker_note(slide,
        "Post-harvest losses of 30–40% represent both a food security failure and an economic "
        "loss to farmers. Supply chain modernization recovers this value. Ethiopia's commodity "
        "exchange increased farmer prices 22–30% by creating transparent, competitive markets. "
        "Cold-chain investments in East Africa cut mango losses from 45% to 8%.")


# ─────────────────────────────────────────────
# CASE STUDIES  (Slides 11-13)
# ─────────────────────────────────────────────

def _case_header(slide, country, headline):
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.05), C_DARK_GREEN)
    _textbox(slide, Inches(0.35), Inches(0.08), Inches(10.0), Inches(0.5),
             headline, font_size=26, bold=True, color=C_WHITE)
    _textbox(slide, Inches(0.35), Inches(0.55), Inches(6.0), Inches(0.42),
             f"Case Study | {country}", font_size=13, color=C_WARM_AMBER)
    _textbox(slide, Inches(10.5), Inches(0.25), Inches(2.7), Inches(0.5),
             "Case Studies", font_size=11, color=C_WARM_AMBER, align=PP_ALIGN.RIGHT)


def add_case_study_1():
    """Slide 11 – Ethiopia's Agricultural Transformation Agency"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _case_header(slide, "Ethiopia", "Ethiopia: From Food Aid to Export Powerhouse")

    _rect(slide, Inches(0.35), Inches(1.15), Inches(7.5), Inches(5.9), C_LIGHT_GREEN)
    _rect(slide, Inches(8.15), Inches(1.15), Inches(4.8), Inches(5.9), RGBColor(255, 243, 224))

    _textbox(slide, Inches(0.5), Inches(1.2), Inches(7.2), Inches(0.5),
             "Context & Intervention", font_size=16, bold=True, color=C_DARK_GREEN)

    context_bullets = [
        "ATA (Agricultural Transformation Agency) established 2010",
        "Targeted 4 priority crops: wheat, maize, teff, coffee",
        "National soil mapping covering 41,000 km²",
        "Bundled input credit + agronomic advisory + market linkage",
        "1 million+ farmers enrolled in extension programme",
        "Ethiopia Commodity Exchange launched for transparent pricing",
        "Irrigation expansion: 350,000 ha new irrigated farmland",
    ]
    _bullet_list(slide, context_bullets, Inches(0.5), Inches(1.8),
                 Inches(7.2), Inches(5.0), font_size=13)

    _textbox(slide, Inches(8.3), Inches(1.2), Inches(4.5), Inches(0.5),
             "Impact Metrics", font_size=16, bold=True, color=C_EARTH_BROWN)

    metrics = [
        ("Wheat yield", "2.1 → 3.8 t/ha", "+81%"),
        ("Maize yield", "2.9 → 5.5 t/ha", "+90%"),
        ("Coffee exports", "$800M → $1.4B", "+75%"),
        ("Farmer income", "+45% avg. real", "2010–2022"),
        ("Food aid need", "Reduced 60%", "since 2015"),
    ]
    for i, (label, value, delta) in enumerate(metrics):
        y = Inches(1.75) + i * Inches(1.0)
        _rect(slide, Inches(8.3), y, Inches(4.5), Inches(0.85), C_WHITE)
        _textbox(slide, Inches(8.45), y + Inches(0.05), Inches(2.2), Inches(0.38),
                 label, font_size=12, color=C_DARK_GRAY)
        _textbox(slide, Inches(8.45), y + Inches(0.42), Inches(2.5), Inches(0.38),
                 value, font_size=14, bold=True, color=C_DARK_GREEN)
        _textbox(slide, Inches(11.2), y + Inches(0.2), Inches(1.5), Inches(0.45),
                 delta, font_size=14, bold=True, color=C_MID_GREEN, align=PP_ALIGN.RIGHT)

    _add_speaker_note(slide,
        "Ethiopia is one of the most dramatic agricultural transformation stories in Sub-Saharan "
        "Africa. The ATA used a systems approach – soil testing, bundled services, and market "
        "reform simultaneously. Key lesson: transformation requires political commitment at the "
        "highest level (PM Meles personally championed this). The ECX brought price transparency "
        "that broke trader monopolies, directly benefiting 15 million farmers.")


def add_case_study_2():
    """Slide 12 – Netherlands precision agriculture"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _case_header(slide, "Netherlands", "Netherlands: World's #2 Food Exporter from a Tiny Country")

    _rect(slide, Inches(0.35), Inches(1.15), Inches(7.5), Inches(5.9), C_LIGHT_GREEN)
    _rect(slide, Inches(8.15), Inches(1.15), Inches(4.8), Inches(5.9), RGBColor(255, 243, 224))

    _textbox(slide, Inches(0.5), Inches(1.2), Inches(7.2), Inches(0.5),
             "How They Did It", font_size=16, bold=True, color=C_DARK_GREEN)

    context_bullets = [
        "Area: 41,000 km² — smaller than West Virginia",
        "Wageningen University-led precision farming research",
        "Greenhouse technology: 75x average global yield/m²",
        "85% reduction in pesticide use via integrated systems",
        "Water recycling in glasshouses: 90% reuse rate",
        "Full digital traceability from seed to shelf",
        "'Food Valley' cluster: 14,000 food industry companies",
    ]
    _bullet_list(slide, context_bullets, Inches(0.5), Inches(1.8),
                 Inches(7.2), Inches(5.0), font_size=13)

    _textbox(slide, Inches(8.3), Inches(1.2), Inches(4.5), Inches(0.5),
             "Impact Metrics", font_size=16, bold=True, color=C_EARTH_BROWN)

    metrics = [
        ("Food exports", "$105 billion/year", "2nd globally"),
        ("Water usage", "−90% vs. global avg", "per kg produce"),
        ("Pesticides", "−85% since 1985", "sustained"),
        ("Labour productivity", "5x EU average", "per farm"),
        ("R&D investment", "€1.2B/year", "Wageningen"),
    ]
    for i, (label, value, delta) in enumerate(metrics):
        y = Inches(1.75) + i * Inches(1.0)
        _rect(slide, Inches(8.3), y, Inches(4.5), Inches(0.85), C_WHITE)
        _textbox(slide, Inches(8.45), y + Inches(0.05), Inches(2.2), Inches(0.38),
                 label, font_size=12, color=C_DARK_GRAY)
        _textbox(slide, Inches(8.45), y + Inches(0.42), Inches(2.5), Inches(0.38),
                 value, font_size=14, bold=True, color=C_DARK_GREEN)
        _textbox(slide, Inches(11.2), y + Inches(0.2), Inches(1.5), Inches(0.45),
                 delta, font_size=12, color=C_MID_GRAY, align=PP_ALIGN.RIGHT)

    _add_speaker_note(slide,
        "The Netherlands proves that scale of land is not the determinant – technology, R&D, and "
        "knowledge intensity are. Wageningen University is ranked #1 agricultural university "
        "globally and has been the backbone of this transformation since the 1960s. The key "
        "lesson for developing nations: invest in agricultural education and R&D first. "
        "Food Valley is now being replicated in Kenya, Vietnam, and Brazil.")


def add_case_study_3():
    """Slide 13 – India's digital agriculture"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _case_header(slide, "India", "India: Digital Agriculture at Scale")

    _rect(slide, Inches(0.35), Inches(1.15), Inches(7.5), Inches(5.9), C_LIGHT_GREEN)
    _rect(slide, Inches(8.15), Inches(1.15), Inches(4.8), Inches(5.9), RGBColor(255, 243, 224))

    _textbox(slide, Inches(0.5), Inches(1.2), Inches(7.2), Inches(0.5),
             "Key Programmes & Interventions", font_size=16, bold=True, color=C_DARK_GREEN)

    context_bullets = [
        "PM-KISAN: Direct income transfer to 110 million farmers",
        "eNAM: National electronic agriculture marketplace",
        "Soil Health Card scheme: 220 million cards issued",
        "Pradhan Mantri Fasal Bima: Crop insurance 60M farmers",
        "AGRISTACK: Digital farmer identity & land records",
        "KCC (Kisan Credit Card): 70 million credit accounts",
        "Drone policy: 100,000 drones for agriculture by 2025",
    ]
    _bullet_list(slide, context_bullets, Inches(0.5), Inches(1.8),
                 Inches(7.2), Inches(5.0), font_size=13)

    _textbox(slide, Inches(8.3), Inches(1.2), Inches(4.5), Inches(0.5),
             "Impact Metrics", font_size=16, bold=True, color=C_EARTH_BROWN)

    metrics = [
        ("eNAM traders", "1.77 million", "+50% price realised"),
        ("Soil card uptake", "220 million", "fertiliser −10%"),
        ("Crop insurance", "60M farmers", "60% risk coverage"),
        ("Digital credit", "70M KCC", "₹6.6T outstanding"),
        ("Food production", "330M tonnes", "record 2023–24"),
    ]
    for i, (label, value, delta) in enumerate(metrics):
        y = Inches(1.75) + i * Inches(1.0)
        _rect(slide, Inches(8.3), y, Inches(4.5), Inches(0.85), C_WHITE)
        _textbox(slide, Inches(8.45), y + Inches(0.05), Inches(2.2), Inches(0.38),
                 label, font_size=12, color=C_DARK_GRAY)
        _textbox(slide, Inches(8.45), y + Inches(0.42), Inches(2.5), Inches(0.38),
                 value, font_size=14, bold=True, color=C_DARK_GREEN)
        _textbox(slide, Inches(11.2), y + Inches(0.2), Inches(1.5), Inches(0.45),
                 delta, font_size=11, color=C_MID_GRAY, align=PP_ALIGN.RIGHT)

    _add_speaker_note(slide,
        "India shows that transformation at scale – 150 million farm households – is possible "
        "with digital-first policy design. AGRISTACK creates a unified digital identity that "
        "links land records, credit history, and insurance, removing the need for physical "
        "document visits. eNAM broke the 'commission agent monopoly' that had suppressed farmer "
        "prices for decades. Key challenge: last-mile internet connectivity for rural areas.")


# ─────────────────────────────────────────────
# BENEFITS & IMPACT  (Slides 14-16)
# ─────────────────────────────────────────────

def add_benefits_productivity():
    """Slide 14 – Increased productivity"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _content_header(slide, "Benefit 1: Increased Productivity", "Benefits & Impact")

    # Bar chart: yield improvements
    chart_data = ChartData()
    chart_data.categories = ["Wheat", "Maize", "Rice", "Vegetables", "Fruits"]
    chart_data.add_series("Traditional (t/ha)", (3.2, 4.1, 3.8, 12.0, 8.5))
    chart_data.add_series("Transformed (t/ha)", (5.8, 8.5, 6.9, 28.0, 18.0))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.3), Inches(1.15), Inches(7.5), Inches(5.8), chart_data
    )
    chart = chart_frame.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Yield Comparison: Traditional vs. Transformed Farming"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = C_MID_GRAY
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = C_MID_GREEN

    _textbox(slide, Inches(8.1), Inches(1.15), Inches(5.0), Inches(0.5),
             "Key Productivity Gains", font_size=16, bold=True, color=C_DARK_GREEN)

    gains = [
        "Precision inputs increase yields 20–80% depending on crop",
        "Reduced waste: post-harvest losses cut from 40% to <10%",
        "Year-round production via irrigation & protected cultivation",
        "Better varieties: GM/hybrid seeds 30–50% higher yield potential",
        "Optimised spacing & density via digital farm planning",
        "Real-time crop monitoring prevents 15–25% yield loss from disease",
    ]
    _bullet_list(slide, gains, Inches(8.1), Inches(1.75), Inches(5.0), Inches(5.0), font_size=14)

    _add_speaker_note(slide,
        "The chart shows realistic yield improvements achievable with transformation. These are "
        "not theoretical maximums but documented field results from pilot programmes. The key "
        "insight is that productivity growth reduces pressure to expand agricultural land, "
        "protecting forests and biodiversity. Doubling yields on existing land is more "
        "achievable – and more sustainable – than expanding the agricultural frontier.")


def add_benefits_environmental():
    """Slide 15 – Environmental benefits"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _content_header(slide, "Benefit 2: Environmental Benefits", "Benefits & Impact")

    benefits = [
        ("🌍  Carbon Sequestration",
         "Regenerative practices can sequester 1–4 t CO₂/ha/year;\n"
         "global agri-land has potential to offset 30% of current GHG"),
        ("💧  Water Conservation",
         "Smart irrigation and water harvesting reduce agricultural\n"
         "water use 40–60% while maintaining or improving yields"),
        ("🌿  Biodiversity Restoration",
         "Agroforestry and hedgerow integration increases on-farm\n"
         "bird species 30%, insect pollinator populations up 45%"),
        ("🪨  Soil Health Recovery",
         "Conservation tillage + cover crops rebuild 0.5–1% soil\n"
         "organic matter per year, reversing decades of degradation"),
        ("🏭  Reduced Chemical Pollution",
         "IPM and precision application cut pesticide/herbicide\n"
         "runoff 60–80%, restoring aquatic ecosystems"),
        ("🌲  Deforestation Prevention",
         "Higher yields on existing land reduce pressure to clear\n"
         "new forest – estimated 100M ha saved by 2030 (CGIAR)"),
    ]

    for i, (title, body) in enumerate(benefits):
        col = i % 2
        row = i // 2
        x = Inches(0.35) + col * Inches(6.55)
        y = Inches(1.2) + row * Inches(1.9)
        _rect(slide, x, y, Inches(6.2), Inches(1.8),
              C_LIGHT_GREEN if col == 0 else RGBColor(255, 243, 224))
        _textbox(slide, x + Inches(0.15), y + Inches(0.08),
                 Inches(5.9), Inches(0.45),
                 title, font_size=14, bold=True, color=C_DARK_GREEN)
        _textbox(slide, x + Inches(0.15), y + Inches(0.55),
                 Inches(5.9), Inches(1.1),
                 body, font_size=12, color=C_DARK_GRAY)

    _add_speaker_note(slide,
        "The environmental co-benefits of agricultural transformation are enormous and often "
        "undervalued in economic analyses. Healthy soil stores more carbon than all the world's "
        "forests combined. The Bonn Challenge and 4 per 1000 initiative are building momentum "
        "around soil carbon. Biodiversity co-benefits matter for food security too – 75% of "
        "crops depend on insect pollination. Frame environment not as a constraint but as "
        "an asset to be maintained.")


def add_benefits_economic_social():
    """Slide 16 – Economic & social impact"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _content_header(slide, "Benefit 3: Economic & Social Impact", "Benefits & Impact")

    col_w = Inches(5.9)
    left1 = Inches(0.4)
    left2 = Inches(6.9)

    _rect(slide, left1, Inches(1.15), col_w, Inches(0.5), C_DARK_GREEN)
    _textbox(slide, left1, Inches(1.15), col_w, Inches(0.5),
             "  Economic Advantages", font_size=16, bold=True, color=C_WHITE)
    _rect(slide, left1, Inches(1.65), col_w, Inches(5.5), C_LIGHT_GREEN)

    econ = [
        "Farmer income doubles to triples with market linkage",
        "Rural GDP growth: +1.5–2% per 10% yield improvement",
        "$1 ag investment creates $3–5 of wider economic activity",
        "Agri-food processing creates 10x more jobs than raw export",
        "Export diversification: reduces FX vulnerability",
        "Input-cost savings: 20–40% from precision application",
        "Reduced food import bills: national fiscal savings",
    ]
    _bullet_list(slide, econ, left1 + Inches(0.15), Inches(1.75),
                 col_w - Inches(0.3), Inches(5.2), font_size=14)

    _rect(slide, left2, Inches(1.15), col_w, Inches(0.5), C_EARTH_BROWN)
    _textbox(slide, left2, Inches(1.15), col_w, Inches(0.5),
             "  Social Impact", font_size=16, bold=True, color=C_WHITE)
    _rect(slide, left2, Inches(1.65), col_w, Inches(5.5), RGBColor(255, 243, 224))

    social = [
        "Food security: 100–150 million lifted from hunger",
        "Gender equity: digital tools empower women farmers",
        "Youth employment: agri-tech creates 'cool' rural jobs",
        "Nutrition: diversified crops improve dietary quality",
        "Education: higher incomes fund children's schooling",
        "Healthcare: safer pesticide practices, less illness",
        "Reduced rural-urban migration and urban pressure",
    ]
    _bullet_list(slide, social, left2 + Inches(0.15), Inches(1.75),
                 col_w - Inches(0.3), Inches(5.2), font_size=14)

    _add_speaker_note(slide,
        "The social returns to agricultural transformation may exceed the economic returns. "
        "Evidence from sub-Saharan Africa shows that a 10% increase in agricultural GDP reduces "
        "poverty by 7% — double the effect of equivalent non-agricultural growth (World Bank). "
        "Women particularly benefit: when women control income, 90% is invested in family "
        "nutrition and education vs. 30–40% for men. Digital tools that give women direct "
        "access to payments bypass traditional gatekeepers.")


# ─────────────────────────────────────────────
# CHALLENGES & SOLUTIONS  (Slides 17-18)
# ─────────────────────────────────────────────

def add_challenges():
    """Slide 17 – Implementation barriers"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _content_header(slide, "Challenges: Implementation Barriers", "Challenges & Solutions")

    barriers = [
        ("Financial",
         ["High upfront technology cost", "Limited rural credit access",
          "$170B global smallholder finance gap", "Short loan tenors vs. long ag cycles"]),
        ("Capacity",
         ["Low digital literacy among farmers", "Shortage of extension workers",
          "Inadequate agricultural education", "Brain drain from rural areas"]),
        ("Infrastructure",
         ["Poor rural roads and connectivity", "Lack of electricity in remote areas",
          "Inadequate cold-chain and storage", "Unreliable water supply"]),
        ("Policy & Governance",
         ["Fragmented policy frameworks", "Weak land tenure and property rights",
          "Bureaucratic subsidy systems", "Political short-termism (election cycles)"]),
        ("Social & Cultural",
         ["Farmer risk-aversion to new methods", "Gender barriers to technology access",
          "Low trust in private-sector partners", "Language and literacy barriers"]),
        ("Climate & Environment",
         ["Increasing weather unpredictability", "Pest/disease range expansion",
          "Degraded natural resource base", "Competing land use pressures"]),
    ]

    card_w = Inches(4.0)
    card_h = Inches(2.65)
    positions = [
        (Inches(0.3),  Inches(1.15)),
        (Inches(4.65), Inches(1.15)),
        (Inches(9.0),  Inches(1.15)),
        (Inches(0.3),  Inches(4.05)),
        (Inches(4.65), Inches(4.05)),
        (Inches(9.0),  Inches(4.05)),
    ]

    for (x, y), (heading, items) in zip(positions, barriers):
        _rect(slide, x, y, card_w, card_h, C_LIGHT_GREEN)
        _textbox(slide, x + Inches(0.15), y + Inches(0.08),
                 card_w - Inches(0.3), Inches(0.42),
                 heading, font_size=14, bold=True, color=C_DARK_GREEN)
        _bullet_list(slide, items, x + Inches(0.1), y + Inches(0.55),
                     card_w - Inches(0.2), Inches(2.0), font_size=12)

    _add_speaker_note(slide,
        "These six barrier categories interact and compound each other. Financial barriers "
        "prevent investment in infrastructure; poor infrastructure makes technology unviable; "
        "weak governance prevents markets from working; cultural barriers slow adoption. "
        "Any strategy must identify the binding constraints for each context. "
        "Surveys of farmers consistently show that finance is the #1 barrier, followed by "
        "access to reliable markets as the #2 concern. Address these two first.")


def add_solutions():
    """Slide 18 – Solutions"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _content_header(slide, "Solutions: Overcoming the Barriers", "Challenges & Solutions")

    solutions = [
        ("Blended Finance & De-risking",
         "Partial guarantee funds reduce lender risk → lower interest. "
         "Climate bonds, green finance, and donor co-investment leverage domestic capital 5:1"),
        ("Digital Extension at Scale",
         "Mobile video advisory, AI chatbots and WhatsApp farmer groups reach 10x more farmers "
         "per officer; lower cost per farmer contact by 80%"),
        ("Infrastructure Investment Prioritisation",
         "Last-mile solar, rural roads, and shared cold-chain deliver highest returns. "
         "Public-private partnerships accelerate deployment with risk sharing"),
        ("Policy Reform & Incentive Alignment",
         "Redirecting poorly-targeted subsidies to productive investment. "
         "Outcome-based contracts with quarterly dashboards create accountability"),
        ("Farmer Organisation & Aggregation",
         "Cooperatives and producer companies provide scale for collective bargaining, "
         "shared equipment, bulk input purchasing and collective market access"),
        ("Inclusive Technology Design",
         "Voice-based, multilingual, low-bandwidth tools for low-literacy users. "
         "Gender-intentional design: access via women's groups, female field agents"),
    ]

    for i, (title, body) in enumerate(solutions):
        col = i % 2
        row = i // 2
        x = Inches(0.35) + col * Inches(6.55)
        y = Inches(1.2) + row * Inches(1.9)
        _rect(slide, x, y, Inches(6.2), Inches(1.8),
              RGBColor(232, 245, 233) if col == 0 else RGBColor(255, 243, 224))
        _rect(slide, x, y, Inches(0.12), Inches(1.8), C_MID_GREEN)
        _textbox(slide, x + Inches(0.2), y + Inches(0.08),
                 Inches(5.85), Inches(0.45),
                 title, font_size=14, bold=True, color=C_DARK_GREEN)
        _textbox(slide, x + Inches(0.2), y + Inches(0.55),
                 Inches(5.85), Inches(1.15),
                 body, font_size=12, color=C_DARK_GRAY)

    _add_speaker_note(slide,
        "Solutions must match the specific barrier profile of each context. In fragile states, "
        "governance reform is the prerequisite. In middle-income countries, finance unlocking "
        "is often the binding constraint. Blended finance instruments (first-loss guarantees, "
        "concessional tranches) have been proven in Kenya, Ghana, and Bangladesh to unlock "
        "10x commercial capital for previously unbankable smallholders.")


# ─────────────────────────────────────────────
# FUTURE OUTLOOK  (Slides 19-20)
# ─────────────────────────────────────────────

def add_future_trends():
    """Slide 19 – Emerging trends"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _content_header(slide, "Future Outlook: Emerging Trends", "Future Outlook")

    trends = [
        ("🤖  AI & Generative AI",
         "GPT-powered crop advisory, autonomous farm robots, AI-designed crop varieties "
         "customised for local microclimates — commercialising 2025–2027"),
        ("🧬  Biotech & Gene Editing",
         "CRISPR crops with drought tolerance, disease resistance and enhanced nutrition "
         "pipeline; 3–5x faster variety development than traditional breeding"),
        ("🌾  Vertical & Controlled Environment",
         "90% less water, zero pesticides, year-round production. Cost parity with field "
         "crops for leafy greens by 2027; protein crops by 2030"),
        ("🔗  Agri-Blockchain & Tokenisation",
         "Carbon credit markets on-chain, tokenised land rights, DeFi agricultural lending "
         "reaching unbanked farmers — $2B market by 2026"),
        ("🦠  Microbiome & Biostimulants",
         "Precision soil biology interventions replacing synthetic fertilisers, with "
         "equivalent yield at 30–50% lower cost; market $18B by 2028"),
        ("🛰️  Space-based Agriculture",
         "Hyperspectral satellites with 3-metre resolution, daily crop monitoring at $0.001/ha; "
         "combined with AI models for sub-field management recommendations"),
    ]

    for i, (title, body) in enumerate(trends):
        col = i % 2
        row = i // 2
        x = Inches(0.35) + col * Inches(6.55)
        y = Inches(1.2) + row * Inches(1.9)
        _rect(slide, x, y, Inches(6.2), Inches(1.8),
              C_LIGHT_GREEN if (col + row) % 2 == 0 else RGBColor(255, 243, 224))
        _textbox(slide, x + Inches(0.15), y + Inches(0.08),
                 Inches(5.9), Inches(0.45),
                 title, font_size=14, bold=True, color=C_DARK_GREEN)
        _textbox(slide, x + Inches(0.15), y + Inches(0.55),
                 Inches(5.9), Inches(1.15),
                 body, font_size=12, color=C_DARK_GRAY)

    _add_speaker_note(slide,
        "These trends are not speculative – all have live pilots or early commercial deployments "
        "in 2024–2025. The convergence of AI, biotech, and digital finance creates exponential "
        "possibilities. Governments and investors that position early will capture disproportionate "
        "benefits. Key challenge: ensuring these technologies are accessible to smallholders, "
        "not only commercial farms.")


def add_opportunities():
    """Slide 20 – Opportunities ahead"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_CREAM)
    _content_header(slide, "Opportunities Ahead: The Next Decade", "Future Outlook")

    _textbox(slide, Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.5),
             "Strategic windows of opportunity for governments, investors, and farmers",
             font_size=16, italic=True, color=C_EARTH_BROWN)

    opportunities = [
        ("Market Growth",
         "$8.4 trillion global food market by 2030 — organic, premium, and functional food "
         "segments growing 2–3x faster than conventional"),
        ("Climate Finance",
         "$100+ billion/year climate finance available; agriculture increasingly eligible "
         "through carbon markets, REDD+, and green bonds"),
        ("Youth & Jobs",
         "600 million new jobs needed by 2030 — agri-tech, food processing, and cold-chain "
         "logistics offer high-quality rural employment"),
        ("Genomic Revolution",
         "10,000+ new crop varieties in development pipelines; gene-edited varieties will "
         "reach farmers 5–10x faster than previously possible"),
        ("Data Economy",
         "Agricultural data is the new oil — satellite, weather, soil, and market data "
         "services growing to $7B market; farmers can monetise their own data"),
        ("South-South Learning",
         "Countries like Brazil, India, Ethiopia now exporting transformation expertise; "
         "peer learning networks reducing costs and accelerating adoption"),
    ]

    for i, (title, body) in enumerate(opportunities):
        col = i % 2
        row = i // 2
        x = Inches(0.35) + col * Inches(6.55)
        y = Inches(1.75) + row * Inches(1.75)
        _rect(slide, x, y, Inches(0.12), Inches(1.65), C_WARM_AMBER)
        _rect(slide, x + Inches(0.12), y, Inches(6.08), Inches(1.65),
              RGBColor(250, 245, 235))
        _textbox(slide, x + Inches(0.25), y + Inches(0.08),
                 Inches(5.8), Inches(0.45),
                 title, font_size=14, bold=True, color=C_DARK_GREEN)
        _textbox(slide, x + Inches(0.25), y + Inches(0.55),
                 Inches(5.8), Inches(1.0),
                 body, font_size=12, color=C_DARK_GRAY)

    _add_speaker_note(slide,
        "The next decade will determine the trajectory of global food systems for a generation. "
        "Countries that invest now will capture market share, climate finance, and talent. "
        "Those that delay will face rising import bills, climate shocks, and political instability "
        "from food price crises. The window for proactive transformation is narrowing – "
        "climate change is compressing the timeline for action.")


# ─────────────────────────────────────────────
# CALL TO ACTION  (Slide 21)
# ─────────────────────────────────────────────

def add_call_to_action():
    """Slide 21 – Call to action"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_DARK_GREEN)

    # decorative stripes
    _rect(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7), C_MID_GREEN)
    _rect(slide, Inches(0), Inches(6.6), SLIDE_W, Inches(0.12), C_WARM_AMBER)

    _textbox(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
             "Call to Action: Your Next Steps", font_size=36, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)

    _textbox(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.5),
             "Transformation begins with a decision. Here is what each stakeholder can do TODAY:",
             font_size=16, italic=True, color=C_WARM_AMBER, align=PP_ALIGN.CENTER)

    cta_items = [
        ("🏛  GOVERNMENTS",
         ["Commit to a national agricultural transformation strategy with clear 2030 targets",
          "Redirect 30% of input subsidies to productive infrastructure investment",
          "Establish a dedicated Agri-Tech regulatory sandbox for rapid innovation testing",
          "Pass digital land rights legislation to unlock smallholder collateral"]),
        ("🏢  PRIVATE SECTOR",
         ["Scale proven agri-tech solutions in underserved markets",
          "Adopt supply-chain transparency and pay a sustainability premium to farmers",
          "Partner with development finance for blended-finance rural credit products",
          "Invest in 'food-tech for development' R&D, not only high-income markets"]),
        ("🌾  FARMER ORGANISATIONS",
         ["Organise into producer companies to achieve negotiating scale",
          "Adopt digital record-keeping to build credit history and insurance eligibility",
          "Connect to early-warning and advisory SMS/digital platforms",
          "Advocate for policy reforms that remove subsidy-capture by non-farmers"]),
        ("🎓  RESEARCHERS & ACADEMIA",
         ["Prioritise last-mile technology adaptation for smallholder contexts",
          "Publish open-access data and tools for agricultural transformation",
          "Engage farmers as co-creators of technology, not just end-users",
          "Build South-South research collaboration networks"]),
    ]

    col_w = Inches(6.0)
    for i, (actor, actions) in enumerate(cta_items):
        col = i % 2
        row = i // 2
        x = Inches(0.4) + col * (col_w + Inches(0.5))
        y = Inches(1.75) + row * Inches(2.4)
        _rect(slide, x, y, col_w, Inches(2.3), RGBColor(39, 97, 39))
        _textbox(slide, x + Inches(0.15), y + Inches(0.08),
                 col_w - Inches(0.3), Inches(0.45),
                 actor, font_size=14, bold=True, color=C_WARM_AMBER)
        _bullet_list(slide, actions, x + Inches(0.1), y + Inches(0.55),
                     col_w - Inches(0.2), Inches(1.65),
                     font_size=11, bullet="• ", color=C_WHITE)

    _add_speaker_note(slide,
        "End the presentation with a challenge to each stakeholder group. This slide is designed "
        "to be left on screen during Q&A so the audience has something actionable to take away. "
        "Emphasise: the technologies exist, the evidence is in, the financing models are proven. "
        "What is missing is collective will and coordinated action. Agricultural transformation "
        "is not a question of 'if' but 'when' — and early movers gain the most.")


# ─────────────────────────────────────────────
# CONCLUSION  (Slide 22)
# ─────────────────────────────────────────────

def add_conclusion():
    """Slide 22 – Conclusion & Q&A"""
    slide = prs.slides.add_slide(BLANK)
    _solid_bg(slide, C_DARK_GREEN)

    _rect(slide, Inches(0), Inches(5.8), SLIDE_W, Inches(1.7), C_MID_GREEN)
    _rect(slide, Inches(0), Inches(5.65), SLIDE_W, Inches(0.12), C_WARM_AMBER)

    _textbox(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
             "Conclusion", font_size=40, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)

    key_points = [
        "Agricultural transformation is not optional — it is the defining challenge of our era",
        "Technology, sustainability, and smart policy must work together as a system",
        "Evidence from Ethiopia, Netherlands, India and beyond shows it is achievable",
        "Benefits span productivity, environment, economy, and social equity",
        "The window for action is open — but narrowing with each passing year",
    ]

    for i, point in enumerate(key_points):
        _rect(slide, Inches(0.5), Inches(1.15) + i * Inches(0.82),
              Inches(12.3), Inches(0.74), RGBColor(39, 97, 39))
        _textbox(slide, Inches(0.7), Inches(1.22) + i * Inches(0.82),
                 Inches(12.0), Inches(0.6),
                 f"{'0' if i < 9 else ''}{i+1}.  {point}",
                 font_size=15, color=C_WHITE)

    _textbox(slide, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.55),
             "Thank You  |  Q & A  |  Contact: agri-transformation@example.org",
             font_size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    _textbox(slide, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.6),
             "Slides available for download · References & data sources in appendix",
             font_size=13, color=C_LIGHT_GREEN, align=PP_ALIGN.CENTER)

    _add_speaker_note(slide,
        "Summarise the five key takeaways. Invite questions — prepare for these common ones: "
        "(1) How much does transformation cost? (Answer: less than continued subsidising of "
        "the status quo — ROI is 8–15x over 10 years.) "
        "(2) What's the single most important intervention? (Answer: depends on context, but "
        "irrigation + market linkage is the most consistently high-impact combination.) "
        "(3) How long does it take? (Answer: 5–7 years for measurable productivity gains; "
        "15–20 years for full structural transformation.) Thank the audience for their attention.")


# ══════════════════════════════════════════════════════════════════════════════
# Build all slides
# ══════════════════════════════════════════════════════════════════════════════

add_title_slide()           # 1
add_toc_slide()             # 2

add_intro_what()            # 3
add_intro_why()             # 4
add_intro_challenges()      # 5

add_tech_adoption()         # 6
add_sustainable_farming()   # 7
add_climate_smart()         # 8
add_digital_tools()         # 9
add_supply_chain()          # 10

add_case_study_1()          # 11
add_case_study_2()          # 12
add_case_study_3()          # 13

add_benefits_productivity() # 14
add_benefits_environmental()# 15
add_benefits_economic_social() # 16

add_challenges()            # 17
add_solutions()             # 18

add_future_trends()         # 19
add_opportunities()         # 20

add_call_to_action()        # 21
add_conclusion()            # 22

# ══════════════════════════════════════════════════════════════════════════════
# Save to repository root
# ══════════════════════════════════════════════════════════════════════════════

repo_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
output_path = os.path.join(repo_root, "Agricultural_Transformation_Presentation.pptx")
prs.save(output_path)
print(f"✓ PowerPoint created: {output_path}")
print(f"✓ Total slides: {len(prs.slides)}")
